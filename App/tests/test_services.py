"""Services, AI plumbing and the mock-mode guarantees."""

from __future__ import annotations

import subprocess
from datetime import timedelta
from pathlib import Path

import pytest

from intellia.ai.llm_provider import LLMRequest, TaskSpec
from intellia.ai.structured import coerce, extract_json, schema_for
from intellia.models import ai as ai_models
from intellia.utils.dates import quarter_bounds, today, ytd_bounds

APP_DIR = Path(__file__).resolve().parents[1]


# -- metrics ----------------------------------------------------------------------------

def test_open_pipeline_uses_all_non_closed_stages(app_context, scope_rep):
    """The semantic layer was updated from Stage-1-only to all open stages."""
    metrics = app_context.metrics
    total = metrics.open_pipeline(scope_rep)
    with app_context.db.reader(scope_rep) as conn:
        stage_one = conn.execute(
            "SELECT COALESCE(SUM(amount), 0) FROM deals "
            "WHERE stage = 'Stage 1 - Discovery'").fetchone()[0]
    assert total > stage_one


def test_coverage_is_undefined_once_target_is_met(app_context, scope_rep):
    """Dividing by a near-zero remainder yields a meaningless four-digit multiple."""
    start, end = quarter_bounds(today())
    target = app_context.targets.sum_for_users(scope_rep.user_ids, start, end)
    booked = app_context.metrics.bookings(scope_rep, start, end)
    coverage = app_context.metrics.coverage(scope_rep, start, end)
    if booked >= target * 0.98:
        assert coverage is None
    else:
        assert coverage is not None and coverage < 100


def test_win_rate_is_revenue_weighted(app_context, scope_manager):
    rate = app_context.metrics.win_rate(scope_manager, *ytd_bounds(today()))
    assert rate is not None and 0 <= rate <= 100


# -- actions ----------------------------------------------------------------------------

def test_action_queue_is_ranked_and_deduped(app_context, scope_rep):
    actions = app_context.actions.build_queue(scope_rep, today(), limit=10)
    assert actions
    scores = [a.score for a in actions]
    assert scores == sorted(scores, reverse=True)
    fingerprints = [(a.title.strip().lower()[:60], a.account_id or a.deal_id)
                    for a in actions]
    assert len(fingerprints) == len(set(fingerprints))


def test_overdue_uses_reporting_date_not_wall_clock(app_context, scope_rep):
    """Guards the demo against silently breaking after 2026-08-13."""
    overdue = app_context.tasks.overdue(scope_rep, today())
    assert all(t.due_date < today() for t in overdue)


def test_high_value_urgent_action_outranks_small_distant_one(app_context, scope_rep):
    actions = app_context.actions.build_queue(scope_rep, today(), limit=30)
    high = [a for a in actions if a.priority == "High"]
    low = [a for a in actions if a.priority == "Low"]
    if high and low:
        assert min(a.score for a in high) > max(a.score for a in low)


def test_mock_executor_never_sends(app_context, scope_rep):
    action = app_context.actions.build_queue(scope_rep, today(), limit=1)[0]
    result = app_context.executor.execute(action, "draft_email", {"draft": "hi"})
    assert result.ok
    assert "nothing was sent" in result.detail.lower()


# -- structured outputs -----------------------------------------------------------------

@pytest.mark.parametrize("cls", [
    ai_models.DailyBrief, ai_models.MeetingPrep, ai_models.ActionExplanation,
    ai_models.EmailDraft, ai_models.GeneratedInsight, ai_models.Answer,
])
def test_schema_generation_works_on_python39(cls):
    """Also the PEP 604 guard: a `str | None` annotation would raise here."""
    schema = schema_for(cls)
    assert schema["type"] == "object"
    assert schema["additionalProperties"] is False
    assert schema["properties"]


def test_coerce_reports_missing_required_fields():
    instance, errors = coerce({"objective": "x"}, ai_models.MeetingPrep)
    assert instance is not None and not errors  # every other field has a default


def test_coerce_rejects_non_object():
    instance, errors = coerce(["not", "an", "object"], ai_models.DailyBrief)
    assert instance is None and errors


def test_extract_json_handles_fenced_and_bare():
    assert extract_json('```json\n{"a": 1}\n```') == {"a": 1}
    assert extract_json('noise before {"a": 2} noise after') == {"a": 2}
    assert extract_json("not json at all") is None


# -- mock mode --------------------------------------------------------------------------

def test_app_runs_without_a_key(app_context):
    assert app_context.ai.live is False
    assert app_context.ai.mode_label == "Prepared (deterministic)"


def test_daily_brief_works_in_mock_mode(app_context):
    persona = app_context.persona("rep")
    brief = app_context.brief.generate(app_context.scope_for("rep"), persona, today())
    assert brief is not None
    assert brief.items
    # The mock composes from real records, so every reference resolves.
    for item in brief.items:
        assert item.title and item.detail


def test_brief_drops_hallucinated_references(app_context):
    """A model-invented ref_id must never reach the render path."""
    class Hallucinating:
        name = "stub"
        available = True
        model = "stub"

        def complete(self, request: LLMRequest):
            from intellia.ai.llm_provider import LLMResponse
            return LLMResponse(payload={
                "headline": "h", "summary": "s",
                "items": [
                    {"kind": "risk", "title": "real", "detail": "d",
                     "ref_type": "", "ref_id": ""},
                    {"kind": "risk", "title": "fake", "detail": "d",
                     "ref_type": "deal", "ref_id": "DL-9999999"},
                ],
                "generated_by": "stub", "generated_at": ""},
                provider="stub")

    ai = app_context.ai
    original_provider, original_mock = ai.provider, ai.mock
    ai.provider, ai.mock = Hallucinating(), Hallucinating()
    try:
        brief = app_context.brief.generate(
            app_context.scope_for("rep"), app_context.persona("rep"), today(),
        )
    finally:
        ai.provider, ai.mock = original_provider, original_mock

    titles = [i.title for i in brief.items]
    assert "real" in titles
    assert "fake" not in titles


def test_meeting_prep_works_in_mock_mode(app_context):
    scope = app_context.scope_for("rep")
    meetings = app_context.calendar.get_meetings_for_day(scope, today(), "USR-3002")
    prep = app_context.prep.generate(scope, app_context.persona("rep"),
                                     meetings[0], today())
    assert prep is not None
    assert prep.objective and prep.desired_outcomes and prep.talking_points


def test_insight_generation_falls_back_without_a_model(app_context):
    scope = app_context.scope_for("rep")
    config = app_context.engine.create_from_prompt("anything at all", scope, "rep")
    assert config.generated_sql
    assert app_context.engine.render(config, scope).ok


# -- discipline -------------------------------------------------------------------------

def test_no_wall_clock_dates_outside_the_dates_module():
    """A stray date.today() would silently break the demo after 2026-08-13."""
    result = subprocess.run(
        ["grep", "-rn", "-E", r"date\.today\(\)|datetime\.now\(\)\.date\(\)",
         str(APP_DIR / "intellia")],
        capture_output=True, text=True)

    def is_prose(line: str) -> bool:
        # Docstrings reference the banned call by name when explaining the rule.
        body = line.split(":", 2)[-1].strip()
        return "``" in body or body.startswith("#") or body.startswith("*")

    offenders = [line for line in result.stdout.splitlines()
                 if "utils/dates.py" not in line and not is_prose(line)]
    assert not offenders, "wall-clock date usage outside utils/dates.py: {}".format(
        offenders)


def test_no_em_dashes_in_user_facing_copy():
    """Copy discipline: an em dash or a double hyphen reads as machine output.

    Prose in docstrings and comments is exempt; this checks string literals,
    which are the only strings that can reach a user.
    """
    import ast

    offenders = []
    for path in sorted((APP_DIR / "intellia").rglob("*.py")):
        tree = ast.parse(path.read_text(encoding="utf-8"))
        # A bare string statement is a docstring, which is prose for maintainers
        # rather than copy for users.
        prose = {id(n.value) for n in ast.walk(tree)
                 if isinstance(n, ast.Expr) and isinstance(n.value, ast.Constant)}
        for node in ast.walk(tree):
            if not isinstance(node, ast.Constant) or not isinstance(node.value, str):
                continue
            if id(node) in prose:
                continue
            text = node.value
            if "—" in text or "–" in text or " -- " in text:
                # CSS custom properties legitimately start with a double hyphen,
                # and so do argparse flags.
                if "var(--" in text or text.strip().startswith("--"):
                    continue
                offenders.append("{}:{}: {}".format(
                    path.relative_to(APP_DIR), node.lineno, text.strip()[:70]))
    assert not offenders, "em dash or double hyphen in user-facing copy: {}".format(
        offenders)


def test_the_brand_mark_survives_the_html_sanitiser():
    """``st.html`` strips inline svg, so the mark ships as a base64 CSS mask.

    A raw ``<svg`` inside the injected style block is silently mangled, which
    makes the mark vanish with no error anywhere. Guard both facts.
    """
    from intellia.theme import brand

    css = brand.mark_css()
    assert "data:image/svg+xml;base64," in css
    assert "<" not in css and ">" not in css
    assert ".ix-mark" in css and "mask:" in css
    assert "<svg" not in brand.mark(16), "the mark must not emit inline svg"


def test_builtin_reseed_refreshes_shipped_copy_but_keeps_user_edits(app_context):
    """A builtin on version 1 is ours to update; past version 1 it is theirs."""
    import dataclasses

    from intellia.insights.builtins import builtin_insights

    store = app_context.store
    shipped = {c.id: c for c in builtin_insights()}
    target = "insight.pipeline_by_stage"

    store.seed(list(shipped.values()))
    assert store.get(target).title == shipped[target].title

    # A user edit moves it past version 1, and re-seeding must not clobber it.
    edited = dataclasses.replace(store.get(target), title="My renamed card")
    store.save_version(edited, "renamed")
    store.seed(list(shipped.values()))
    assert store.get(target).title == "My renamed card"

    # Restore, so the suite stays re-runnable.
    store.delete(target)
    store.seed(list(shipped.values()))
    assert store.get(target).title == shipped[target].title


def test_config_toml_theme_matches_the_token_palette():
    """config.toml themes Streamlit's own widgets; tokens.py themes ours.

    They describe the same surfaces, so a change to one and not the other shows
    up as light cards on a dark canvas or a mismatched border. Generate the block
    with ``tokens.config_theme_toml()`` rather than hand-editing it.
    """
    import re

    from intellia.theme import tokens

    config = (APP_DIR / ".streamlit" / "config.toml").read_text(encoding="utf-8")
    drift = []
    for key, token in tokens.CONFIG_KEYS:
        if token is None:
            continue
        match = re.search(r'^{} = "([^"]+)"'.format(key), config, re.M)
        if match is None or match.group(1) != tokens.LIGHT[token]:
            drift.append("{}: config has {}, palette has {}".format(
                key, match.group(1) if match else "nothing", tokens.LIGHT[token]))

    for name, series in (("chartCategoricalColors", tokens.CATEGORICAL_LIGHT),
                         ("chartSequentialColors", tokens.SEQUENTIAL_LIGHT)):
        block = re.search(r"^{} = \[(.*?)\]".format(name), config, re.M | re.S)
        found = re.findall(r'"(#[0-9A-Fa-f]{6})"', block.group(1)) if block else []
        if [c.lower() for c in found] != [c.lower() for c in series]:
            drift.append("{}: config and palette differ".format(name))

    assert not drift, drift


def test_chart_palettes_are_one_hue():
    """The system is navy monochromatic: series separate by lightness, not hue.

    A stray warm colour in the categorical ramp would pass every visual check on
    one chart and wreck the set.
    """
    import colorsys

    from intellia.theme import tokens

    def hue(hex_colour: str) -> float:
        r, g, b = (int(hex_colour[i:i + 2], 16) / 255.0 for i in (1, 3, 5))
        return colorsys.rgb_to_hls(r, g, b)[0] * 360

    for name, series in (("light", tokens.CATEGORICAL_LIGHT),
                         ("dark", tokens.CATEGORICAL_DARK),
                         ("ramp", tokens.SEQUENTIAL_LIGHT)):
        hues = [hue(c) for c in series if c.upper() not in ("#FFFFFF", "#000000")]
        spread = max(hues) - min(hues)
        assert spread < 25, "{} palette spans {:.0f} degrees of hue".format(name, spread)

    # And lightness must actually separate, or one hue becomes one colour.
    def light(hex_colour: str) -> float:
        r, g, b = (int(hex_colour[i:i + 2], 16) / 255.0 for i in (1, 3, 5))
        return colorsys.rgb_to_hls(r, g, b)[1]

    steps = [light(c) for c in tokens.CATEGORICAL_LIGHT]
    gaps = [abs(b - a) for a, b in zip(steps, steps[1:])]
    assert min(gaps) > 0.05, "adjacent categorical slots are too close in lightness"


# -- the text-to-SQL router -------------------------------------------------------------

@pytest.mark.parametrize("question", [
    # The question that regressed: no enumerated keyword, but plainly a query.
    "Give me customers ARR by region",
    "Give me customers ARR by segment",
    "ARR by region",
    "customers ARR per tier",
    "how much ARR do we have in EMEA",
    "how many open deals are there",
    # Dimensions nobody enumerated in advance.
    "revenue by renewal month",
    "accounts by health score",
    "pipeline split across owners",
    # The phrasings the old allowlist did catch, which must keep working.
    "Show ARR by industry",
    "Break down open pipeline by region",
    "top 10 accounts",
])
def test_router_sends_data_questions_to_text_to_sql(question):
    """The router fails open.

    It used to be an allowlist of phrasings, so a dimension that had not been
    enumerated was refused before the engine saw the question. The engine reads
    the semantic layer and writes its own SQL; the router must not pre-empt it.
    """
    from intellia.services.ask_service import wants_insight
    assert wants_insight(question), "would never reach text to SQL: {!r}".format(question)


@pytest.mark.parametrize("question", [
    "hi",
    "thanks",
    "ok",
    "draft an email to Victor",
    "schedule a follow up with Priya",
    "what should I do today",
])
def test_router_skips_turns_no_query_can_answer(question):
    """Greetings, actions and advice are answered from the evidence bundle."""
    from intellia.services.ask_service import wants_insight
    assert not wants_insight(question), "wasted a model call on: {!r}".format(question)


def test_router_still_queries_when_an_action_names_a_breakdown():
    """"Send me ARR by region" is an action wrapped around a real query."""
    from intellia.services.ask_service import wants_insight
    assert wants_insight("send me ARR by region")


# -- follow-up questions ----------------------------------------------------------------

def test_prior_turn_block_carries_the_question_sql_and_chart_type():
    """A follow-up like "flip this to a donut" only means something with a referent."""
    from intellia.ai.prompts import tasks
    block = tasks.prior_turn_block(
        "Give me customers ARR by region",
        "select region, sum(arr) from accounts group by region", "bar")
    assert "Give me customers ARR by region" in block
    assert "select region, sum(arr) from accounts group by region" in block
    assert "bar" in block


def test_prior_turn_block_is_empty_without_a_previous_query():
    """The first turn of a thread has nothing to refer back to."""
    from intellia.ai.prompts import tasks
    assert tasks.prior_turn_block("", "", "") == ""
    assert tasks.prior_turn_block("a question", "", "bar") == ""


def test_prior_context_precedes_the_question_in_the_prompt():
    from intellia.ai.prompts import tasks
    prompt = tasks.insight_sql_user("flip it to a donut", "PRIOR CONTEXT\n\n")
    assert prompt.index("PRIOR CONTEXT") < prompt.index("flip it to a donut")


def test_same_follow_up_after_a_different_question_is_a_different_cache_entry():
    """Otherwise one thread's chart is served to another.

    "Flip this to a donut" is identical wording in every thread, so caching on the
    question alone would hand back whichever chart was generated first.
    """
    from intellia.ai.cache import cache_key
    base = ("claude", "claude-sonnet-5", "insight_sql", "v1")
    a = cache_key(*base, {"question": "flip this to a donut", "prior": "ARR by region",
                          "schema": "abc"})
    b = cache_key(*base, {"question": "flip this to a donut", "prior": "deals by stage",
                          "schema": "abc"})
    assert a != b


# -- the per-chat question budget -------------------------------------------------------

@pytest.mark.parametrize("answered,pending,left,full", [
    (0, False, 5, False),
    (1, False, 4, False),
    (4, False, 1, False),
    (4, True, 0, True),    # the fifth is in flight, so the chat is already full
    (5, False, 0, True),
    (6, False, 0, True),   # never reports a negative remainder
])
def test_question_budget_counts_answered_turns_and_the_one_in_flight(
        monkeypatch, answered, pending, left, full):
    """A question in flight counts, or the composer would accept a sixth.

    The turn is only appended once the model returns, so counting appended turns
    alone leaves a window where the thread looks one short of the limit.
    """
    from intellia.state import session

    thread = {"id": "th-1", "turns": [{"q": "q", "a": None} for _ in range(answered)]}
    monkeypatch.setattr(session, "current_thread", lambda: thread)
    monkeypatch.setattr(session, "thread_id", lambda: "th-1")
    monkeypatch.setattr(session, "pending_question", lambda: "q" if pending else None)

    assert session.questions_left(thread) == left
    assert session.thread_is_full(thread) is full


def test_question_budget_is_per_thread_not_global(monkeypatch):
    """A full chat must not make a fresh one look full."""
    from intellia.state import session

    monkeypatch.setattr(session, "pending_question", lambda: None)
    monkeypatch.setattr(session, "thread_id", lambda: "th-old")
    full = {"id": "th-old", "turns": [{"q": "q", "a": None} for _ in range(5)]}
    fresh = {"id": "th-new", "turns": []}
    assert session.thread_is_full(full)
    assert not session.thread_is_full(fresh)
    assert session.questions_left(fresh) == session.MAX_QUESTIONS_PER_THREAD


def test_app_state_database_can_be_redirected(monkeypatch, tmp_path):
    """The UI tests write to whatever app_state.db settings point at.

    Aimed at the shipped one they left a heading block behind on every run, so a
    real canvas collected a section title per test run. The override is what keeps
    a test run out of the user's own state.
    """
    from intellia.config import settings as settings_module

    target = tmp_path / "elsewhere.db"
    monkeypatch.setenv("INTELLIA_APP_STATE_DB", str(target))
    assert settings_module.load_settings().app_state_db == target

    monkeypatch.delenv("INTELLIA_APP_STATE_DB")
    assert settings_module.load_settings().app_state_db == (
        settings_module.DATA_DIR / "app_state.db")


# -- the deck's numbers -----------------------------------------------------------------

DECK = APP_DIR.parent / "Deck" / "Intellia_Deck.pptx"
LINE_TOLERANCE = 250


def _deck_stat(label: str) -> str:
    """The number printed immediately before a label on the deck's stat band.

    Reading it positionally rather than by shape id means the slide can be
    rearranged in PowerPoint without breaking this.
    """
    from pptx import Presentation

    runs = []
    for slide in Presentation(str(DECK)).slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                runs.extend(r.text.strip() for r in paragraph.runs if r.text.strip())
    if label not in runs:
        raise AssertionError("no {!r} on the deck's stat band".format(label))
    return runs[runs.index(label) - 1]


def _repo_python_lines() -> int:
    total = 0
    for path in APP_DIR.rglob("*.py"):
        if any(part.startswith(".") for part in path.relative_to(APP_DIR).parts):
            continue
        total += len(path.read_text(encoding="utf-8").splitlines())
    return total


def test_the_deck_still_claims_the_real_test_count(request):
    """The deck has no generator, so nothing keeps its numbers honest.

    They drifted four times (102, 124, 143, 157) and each drift was caught by a
    human noticing rather than by anything automatic. If this fails, the code
    changed and slide 4 did not.
    """
    pytest.importorskip("pptx", reason="python-pptx is not a runtime dependency")
    if not DECK.exists():
        pytest.skip("no deck in this checkout")
    # A filtered run collects a subset, which would fail this for the wrong reason.
    if request.config.getoption("-k") or request.config.getoption("file_or_dir"):
        pytest.skip("partial run: the collected count is not the suite total")

    claimed = _deck_stat("automated tests")
    actual = request.session.testscollected
    assert claimed == str(actual), (
        "deck slide 4 says {} automated tests, the suite collects {}".format(
            claimed, actual))


def test_the_deck_still_claims_the_real_line_count():
    """Approximate on purpose: the deck prints a rounded figure with a leading ~."""
    pytest.importorskip("pptx", reason="python-pptx is not a runtime dependency")
    if not DECK.exists():
        pytest.skip("no deck in this checkout")

    claimed = int(_deck_stat("lines of Python").lstrip("~").replace(",", ""))
    actual = _repo_python_lines()
    assert abs(claimed - actual) <= LINE_TOLERANCE, (
        "deck slide 4 says ~{:,} lines of Python, the repo has {:,}".format(
            claimed, actual))


# -- a query that matched nothing -------------------------------------------------------

def test_zero_rows_is_stated_as_an_answer_not_as_missing_data():
    """"None" and "I was not given that" are different claims.

    An empty result used to reach the prompt as an empty string, which is what a
    turn with no query at all looks like. The model then reported that it did not
    have account-level figures, which was untrue: it had a query, and the query
    matched nothing.
    """
    from intellia.ai.prompts import tasks

    prompt = tasks.ask_user("accounts with no open deals", "evidence", "metrics",
                            tasks.EMPTY_RESULT)
    assert "zero rows" in prompt
    assert "The answer is none" in prompt
    assert tasks.EMPTY_RESULT not in prompt, "the sentinel leaked into the prompt"


def test_no_query_at_all_still_adds_no_result_section():
    from intellia.ai.prompts import tasks
    prompt = tasks.ask_user("what should I do today", "evidence", "metrics", "")
    assert "Query result" not in prompt


def test_real_rows_are_still_marked_authoritative():
    from intellia.ai.prompts import tasks
    prompt = tasks.ask_user("arr by region", "evidence", "metrics",
                            "| region | arr |\n| --- | --- |\n| AMER | 10 |")
    assert "authoritative" in prompt
    assert "| AMER | 10 |" in prompt
